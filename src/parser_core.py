# src/parser_core.py
import re
from pathlib import Path
import json
import csv
import warnings

# Librerías de terceros
import pdfplumber
import docx
from bs4 import BeautifulSoup, XMLParsedAsHTMLWarning
import pptx
import openpyxl
from odf import text, teletype
from odf.opendocument import load as load_odt
from striprtf.striprtf import rtf_to_text
from ebooklib import epub, ITEM_DOCUMENT
from markdown import markdown

# --- FUNCIONES DE PROCESAMIENTO DE TEXTO ---

def limpiar_y_estructurar_texto(texto_bruto: str) -> str:
    """
    Toma el texto en bruto y aplica reglas de limpieza y estructuración a Markdown.
    """
    # Patrones de ruido a eliminar (cabeceras, pies de página, etc.)
    noise_patterns = [
        r'copyright', r'todos los derechos reservados', r'aviso legal',
        r'política de privacidad', r'agencia de traducción', r'traducciones profesionales',
        r'ibidem group', r'contacto', r'icono cabecera'
    ]
    
    # Patrones para identificar elementos estructurales
    code_prompt_pattern = re.compile(r'^\s*(?:>>>|\$|#|~|\.\.\.)\s')
    title_pattern = re.compile(r'^\s*[A-Z\s]{5,50}\s*$') # Títulos en mayúsculas (5 a 50 chars)
    list_item_pattern = re.compile(r'^\s*[-*•]\s+|^\s*\d+\.\s+')

    lineas_originales = texto_bruto.split('\n')
    lineas_procesadas = []
    
    in_code_block = False

    for linea in lineas_originales:
        linea_limpia = linea.strip()

        if not linea_limpia or any(re.search(p, linea_limpia, re.IGNORECASE) for p in noise_patterns) or linea_limpia.isdigit():
            continue

        if code_prompt_pattern.match(linea):
            if not in_code_block:
                lineas_procesadas.append("\n```python")
                in_code_block = True
            lineas_procesadas.append(linea)
            continue
        elif in_code_block:
            lineas_procesadas.append("```\n")
            in_code_block = False

        if title_pattern.match(linea_limpia) and len(linea_limpia.split()) < 10:
            lineas_procesadas.append(f"## {linea_limpia}\n")
            continue
        
        if list_item_pattern.match(linea_limpia):
            linea_lista = re.sub(r'^\s*[-*•]', '-', linea_limpia, 1)
            lineas_procesadas.append(linea_lista)
            continue
            
        lineas_procesadas.append(linea_limpia)

    if in_code_block:
        lineas_procesadas.append("```\n")
        
    texto_completo = "\n".join(lineas_procesadas)
    
    texto_completo = re.sub(r'\n{3,}', '\n\n', texto_completo)
    texto_completo = re.sub(r'(\n## .*?)\n+', r'\1\n\n', texto_completo)
    texto_completo = re.sub(r'\n(```)', r'\n\n\1', texto_completo)
    texto_completo = re.sub(r'(```)\n', r'\1\n\n', texto_completo)

    return texto_completo.strip()

# --- CONSTANTE Y FUNCIÓN PRINCIPAL DE EXTRACCIÓN ---

EXTRACTORES = {
    '.pdf': '_extraer_pdf', '.docx': '_extraer_docx', '.txt': '_extraer_txt',
    '.html': '_extraer_html', '.xml': '_extraer_xml', '.pptx': '_extraer_pptx',
    '.xlsx': '_extraer_xlsx', '.odt': '_extraer_odt', '.rtf': '_extraer_rtf',
    '.epub': '_extraer_epub', '.md': '_extraer_md', '.json': '_extraer_json',
    '.csv': '_extraer_csv',
}

def extraer_texto(ruta_archivo: str, progress_callback=None) -> str:
    """
    Toma la ruta de un archivo, detecta su tipo y extrae el texto en bruto.
    """
    ruta = Path(ruta_archivo)
    try:
        if not ruta.is_file():
            return "Error: La ruta no corresponde a un archivo."
        nombre_funcion_extractor = EXTRACTORES.get(ruta.suffix.lower())
        if nombre_funcion_extractor:
            extractor_func = globals()[nombre_funcion_extractor]
            return extractor_func(ruta, progress_callback)
        else:
            return f"Error: Formato de archivo '{ruta.suffix}' no soportado."
    except Exception as e:
        return f"Error al procesar el archivo '{ruta.name}': {e}"

# --- FUNCIONES AUXILIARES DE EXTRACCIÓN (sin cambios) ---
def _extraer_pdf(ruta, cb):
    with pdfplumber.open(ruta) as pdf:
        total_paginas = len(pdf.pages); texto_paginas = []
        for i, page in enumerate(pdf.pages):
            texto_paginas.append(page.extract_text(x_tolerance=1, y_tolerance=1) or "")
            if cb: cb(int(((i + 1) / total_paginas) * 100))
        return "\n".join(texto_paginas)
def _extraer_docx(ruta, cb):
    doc = docx.Document(ruta); texto = "\n".join(para.text for para in doc.paragraphs)
    if cb: cb(100); return texto
def _extraer_txt(ruta, cb):
    texto = ruta.read_text(encoding='utf-8', errors='ignore')
    if cb: cb(100); return texto
def _extraer_html(ruta, cb):
    html_content = ruta.read_text(encoding='utf-8', errors='ignore')
    soup = BeautifulSoup(html_content, 'lxml'); texto = soup.get_text(separator='\n', strip=True)
    if cb: cb(100); return texto
def _extraer_xml(ruta, cb):
    xml_content = ruta.read_text(encoding='utf-8', errors='ignore')
    soup = BeautifulSoup(xml_content, 'lxml-xml'); texto = soup.get_text(separator='\n', strip=True)
    if cb: cb(100); return texto
def _extraer_pptx(ruta, cb):
    prs = pptx.Presentation(ruta); texto = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"): texto.append(shape.text)
    if cb: cb(100); return "\n".join(texto)
def _extraer_xlsx(ruta, cb):
    workbook = openpyxl.load_workbook(ruta, read_only=True); texto = []
    for sheetname in workbook.sheetnames:
        sheet = workbook[sheetname]
        for row in sheet.iter_rows():
            row_text = [str(cell.value) for cell in row if cell.value is not None]
            texto.append(" ".join(row_text))
    if cb: cb(100); return "\n".join(texto)
def _extraer_odt(ruta, cb):
    doc = load_odt(ruta); all_paras = doc.getElementsByType(text.P)
    texto = "\n".join(teletype.extractText(para) for para in all_paras)
    if cb: cb(100); return texto
def _extraer_rtf(ruta, cb):
    rtf_content = ruta.read_text(encoding='ascii', errors='ignore')
    texto = rtf_to_text(rtf_content);
    if cb: cb(100); return texto
def _extraer_epub(ruta, cb):
    warnings.filterwarnings("ignore", category=XMLParsedAsHTMLWarning)
    book = epub.read_epub(ruta); texto = []
    items = list(book.get_items_of_type(ITEM_DOCUMENT)); total_items = len(items)
    for i, item in enumerate(items):
        soup = BeautifulSoup(item.get_content(), 'lxml')
        texto.append(soup.get_text(separator='\n', strip=True))
        if cb: cb(int(((i + 1) / total_items) * 100))
    return "\n".join(texto)
def _extraer_md(ruta, cb):
    md_content = ruta.read_text(encoding='utf-8', errors='ignore')
    html = markdown(md_content); soup = BeautifulSoup(html, 'lxml')
    texto = soup.get_text(separator='\n', strip=True)
    if cb: cb(100); return texto
def _extraer_json(ruta, cb):
    json_data = json.loads(ruta.read_text(encoding='utf-8', errors='ignore'))
    texto = json.dumps(json_data, indent=2, ensure_ascii=False)
    if cb: cb(100); return texto
def _extraer_csv(ruta, cb):
    with ruta.open(mode='r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f); texto = "\n".join(["\t".join(row) for row in reader])
    if cb: cb(100); return texto
